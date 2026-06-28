#!/usr/bin/env python3
"""Generate a clean, professional, editable PPTX for the Gleame technical deep dive.
Design: white bg, near-black text, one restrained deep-blue accent, Arial, strict grid.
Import to Google Slides via Drive -> Open with Google Slides (auto-converts, stays editable).
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ---- palette ----
INK    = RGBColor(0x1A,0x1A,0x1A)
SUB    = RGBColor(0x5B,0x64,0x70)
HAIR   = RGBColor(0xD7,0xDC,0xE3)
PANEL  = RGBColor(0xF4,0xF6,0xF9)
ACCENT = RGBColor(0x1F,0x4E,0x79)   # deep slate blue
ACCENT_SOFT = RGBColor(0xE7,0xED,0xF4)
POS    = RGBColor(0x1E,0x7A,0x46)
NEG    = RGBColor(0xB4,0x23,0x2A)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
MONO   = "Consolas"
FONT   = "Arial"

ML, CW = 0.85, 11.633
RIGHT = ML + CW

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ---- helpers ----
def slide():
    return prs.slides.add_slide(BLANK)

def _run(p, text, size, color=INK, bold=False, italic=False, name=FONT):
    r = p.add_run(); r.text = text
    f = r.font; f.size = Pt(size); f.name = name; f.bold = bold; f.italic = italic
    f.color.rgb = color
    return r

def textbox(s, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    b = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = b.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    return b, tf

def line(s, x1, y1, x2, y2, color=HAIR, w=0.75, dashed=False):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(w); c.shadow.inherit = False
    if dashed:
        ln = c.line._get_or_add_ln()
        ln.append(ln.makeelement(qn('a:prstDash'), {'val': 'dash'}))
    return c

def arrow(s, x1, y1, x2, y2, color=SUB, w=1.25, dashed=False):
    c = line(s, x1, y1, x2, y2, color, w, dashed)
    ln = c.line._get_or_add_ln()
    ln.append(ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    return c

def box(s, l, t, w, h, fill=WHITE, border=HAIR, bw=1.0, rounded=True):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if border is None: shp.line.fill.background()
    else: shp.line.color.rgb = border; shp.line.width = Pt(bw)
    shp.shadow.inherit = False
    if rounded:
        try: shp.adjustments[0] = 0.045
        except Exception: pass
    tf = shp.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08); tf.margin_bottom = Inches(0.08)
    return shp

def labelbox(s, l, t, w, h, label, sub=None, fill=WHITE, border=HAIR,
             lc=INK, lsize=12.5, lbold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    shp = box(s, l, t, w, h, fill=fill, border=border)
    shp.text_frame.vertical_anchor = anchor
    p = shp.text_frame.paragraphs[0]; p.alignment = align
    _run(p, label, lsize, lc, bold=lbold)
    if sub:
        sp = shp.text_frame.add_paragraph(); sp.alignment = align; sp.space_before = Pt(2)
        _run(sp, sub, 9.5, SUB)
    return shp

def header(s, eyebrow, title, idx):
    _, tf = textbox(s, ML, 0.52, CW, 0.32)
    _run(tf.paragraphs[0], eyebrow.upper(), 11, ACCENT, bold=True)
    _, tf2 = textbox(s, ML, 0.84, CW, 0.8)
    _run(tf2.paragraphs[0], title, 25, INK, bold=True)
    rule = box(s, ML, 1.6, 0.62, 0.045, fill=ACCENT, border=None, rounded=False)
    footer(s, idx)

def footer(s, idx):
    _, tf = textbox(s, ML, 7.04, 6, 0.3)
    _run(tf.paragraphs[0], "Gleame  ·  Technical Deep Dive", 9, SUB)
    _, tf2 = textbox(s, RIGHT-1.2, 7.04, 1.2, 0.3)
    tf2.paragraphs[0].alignment = PP_ALIGN.RIGHT
    _run(tf2.paragraphs[0], f"{idx:02d}", 9, SUB)

def card(s, l, t, w, h, title, body, accent=False, body_size=13.5):
    shp = box(s, l, t, w, h, fill=(ACCENT_SOFT if accent else WHITE),
              border=(ACCENT if accent else HAIR), bw=(1.25 if accent else 1.0))
    tf = shp.text_frame; tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]; _run(p, title.upper(), 10.5, ACCENT, bold=True)
    bp = tf.add_paragraph(); bp.space_before = Pt(6); _run(bp, body, body_size, INK)
    return shp

def bullets(s, l, t, w, h, items, size=15, gap=9):
    _, tf = textbox(s, l, t, w, h)
    for k, it in enumerate(items):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        _run(p, "—  ", size, ACCENT, bold=True)
        if isinstance(it, tuple):
            _run(p, it[0], size, INK, bold=True); _run(p, it[1], size, INK)
        else:
            _run(p, it, size, INK)

def row_table(s, l, t, w, rows, colx, fs=13, rowh=0.5, header_row=True):
    """Clean underlined-rows table. colx = fractional x starts (0..1)."""
    y = t
    xs = [l + cx * w for cx in colx] + [l + w]
    line(s, l, y, l + w, y, color=INK, w=1.0)   # top rule
    for ri, row in enumerate(rows):
        is_h = header_row and ri == 0
        for ci, val in enumerate(row):
            cw = xs[ci + 1] - xs[ci]
            _, tf = textbox(s, xs[ci], y + 0.02, cw - 0.1, rowh, anchor=MSO_ANCHOR.MIDDLE)
            p = tf.paragraphs[0]
            if isinstance(val, tuple):
                _run(p, val[0], fs, val[1], bold=val[2])
            else:
                _run(p, val, fs, SUB if is_h else INK, bold=is_h)
        y += rowh
        line(s, l, y, l + w, y, color=(INK if is_h else HAIR), w=(1.0 if is_h else 0.75))
    return y

def notes(s, text):
    s.notes_slide.notes_text_frame.text = text

n = 0
def idx():
    global n; n += 1; return n

# ============================================================== SLIDES

# 1 TITLE
s = slide()
box(s, 0, 0, 0.28, 7.5, fill=ACCENT, border=None, rounded=False)  # thin left spine
_, tf = textbox(s, ML, 2.35, CW, 0.4); _run(tf.paragraphs[0], "TECHNICAL DEEP DIVE", 12, ACCENT, bold=True)
_, tf = textbox(s, ML, 2.75, CW, 1.1); _run(tf.paragraphs[0], "Gleame", 52, INK, bold=True)
_, tf = textbox(s, ML, 3.95, CW, 0.8)
_run(tf.paragraphs[0], "An AI conversion assistant for Shopify beauty stores — virtual try-on and personalized product recommendation.", 18, SUB)
box(s, ML, 5.0, 0.62, 0.045, fill=ACCENT, border=None, rounded=False)
_, tf = textbox(s, ML, 5.2, CW, 0.6)
p = tf.paragraphs[0]; _run(p, "Charlie Gao", 15, INK, bold=True)
sp = tf.add_paragraph(); _run(sp, "Second-round interview · Neo / Foothill Labs", 13, SUB)
notes(s, "Open calmly. One sentence: Gleame is a conversion-optimization system for Shopify beauty stores — the AI try-on is the engine, the product is the funnel, the recommender, and the measurement. State the plan: ~20 minutes, then questions; you will walk through the system, the production results, and the one thing you would fix.")

# 2 AGENDA
s = slide(); header(s, "Overview", "Agenda", idx())
items = [
    ("1 — The system", "Architecture, the cold-start recommender, the image pipeline, security, and measurement."),
    ("2 — Production results", "A live funnel and a Shopify-native A/B test."),
    ("3 — Diagnosis", "Two independent signals resolving to one root cause."),
    ("4 — Roadmap", "The highest-leverage next change."),
]
y = 2.05
for t_, b_ in items:
    card(s, ML, y, CW, 1.0, t_, b_, body_size=14); y += 1.18
notes(s, "Promise the payoff up front so the build section reads as setup, not the whole talk. The thread: the model works; the engineering of interest is everything around it, and the rigor is in the measurement.")

# 3 MOTIVATION
s = slide(); header(s, "Motivation", "Why this exists", idx())
card(s, ML, 2.05, (CW-0.4)/2, 1.7, "Problem 1 — Conversion",
     "Shoppers cannot tell whether a shade suits their skin from a flat product photo. Uncertainty suppresses purchases.", body_size=14.5)
card(s, ML+(CW-0.4)/2+0.4, 2.05, (CW-0.4)/2, 1.7, "Problem 2 — Returns",
     "When they buy and the shade is wrong, it is returned — eroding margin and trust.", body_size=14.5)
b = box(s, ML, 4.1, CW, 1.7, fill=PANEL, border=HAIR)
tf = b.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; _run(p, "Both problems share one cause: the product cannot be experienced before purchase. ", 15.5, INK)
_run(p, "Gleame personalizes with a couple of questions and a selfie, recommends real in-catalog products, and renders each one on the shopper's own face.", 15.5, INK)
notes(s, "Two structural costs, one root cause. Existing options are weak: AR filters do not map to real SKUs; quizzes do not show anything. Gleame's wedge is showing the recommended product on the shopper.")

# 4 PRODUCT OVERVIEW
s = slide(); header(s, "Overview", "What the shopper experiences", idx())
bw_ = 3.2; gap = (CW - 3*bw_)/2; ty = 2.9; bh = 1.5
steps = [("1 · Consult", "A couple of guided questions (e.g. undertone)."),
         ("2 · Selfie", "Photo captured or uploaded; analyzed automatically."),
         ("3 · Recommendations", "Three products, each rendered on the shopper's face.")]
for i_, (t_, b_) in enumerate(steps):
    x = ML + i_*(bw_+gap)
    labelbox(s, x, ty, bw_, bh, t_, sub=b_, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    if i_ < 2:
        arrow(s, x+bw_+0.04, ty+bh/2, x+bw_+gap-0.04, ty+bh/2, color=ACCENT, w=1.5)
_, tf = textbox(s, ML, 4.9, CW, 0.6)
_run(tf.paragraphs[0], "End to end in roughly twenty seconds: from a couple of taps to three shades shown on the shopper.", 14, SUB)
notes(s, "A founder should picture this on their own store. BEFORE THE TALK: consider replacing this with a real 15-second screen recording or three real screenshots — a real before/after is the strongest single frame.")

# 5 SCOPE
s = slide(); header(s, "Scope", "The assistant is five subsystems", idx())
cards = [
    ("Cold-start recommender", "Elicits and infers preferences with no behavioral data."),
    ("Image-generation pipeline", "Multi-model, parallel, fallback-hardened try-ons."),
    ("Deterministic prompt generation", "Funnel-driven, reproducible, merchant-controlled."),
    ("Multi-tenant security boundary", "Domain verification, billing gate, tiered rate limits."),
    ("Analytics & attribution", "Funnel events reconciled to orders via cart token."),
    ("Storefront surface", "Vanilla-JS Shopify theme extension — on-domain, no iframe."),
]
cw3 = (CW-0.8)/3; ch = 1.55
for i_, (t_, b_) in enumerate(cards):
    r, c = divmod(i_, 3)
    x = ML + c*(cw3+0.4); yy = 2.05 + r*(ch+0.35)
    card(s, x, yy, cw3, ch, t_, b_, accent=(i_==5), body_size=13)
notes(s, "Reframe: this is a CRO system, not a try-on widget. The model is one component; these five subsystems are the engineering. The assistant is also the highest-intent surface in the app.")

# 6 FUNNEL FRAME
s = slide(); header(s, "Framing", "The funnel the system optimizes", idx())
stages = ["Open the assistant", "Consultation (a couple of questions)", "Upload a selfie",
          "See personalized recommendations", "Click through and purchase"]
widths = [9.5, 8.6, 7.2, 5.8, 4.4]
yy = 2.1
for st, wd in zip(stages, widths):
    x = ML + (CW - wd)/2
    labelbox(s, x, yy, wd, 0.62, st, fill=PANEL, border=HAIR, align=PP_ALIGN.CENTER, lsize=13)
    yy += 0.78
_, tf = textbox(s, ML, 6.1, CW, 0.6)
_run(tf.paragraphs[0], "Every design decision exists to move shoppers down this funnel — and to measure where they leave.", 14, SUB)
notes(s, "This is the spine. Plant the shape now; we return to it with real numbers. A deliberate brief beat.")

# 7 ARCHITECTURE
s = slide(); header(s, "Architecture", "System architecture", idx())
# groups
def group(s, l, t, w, h, label, col=SUB):
    box(s, l, t, w, h, fill=None, border=HAIR, bw=1.0)
    _, tf = textbox(s, l+0.1, t-0.34, w, 0.3); _run(tf.paragraphs[0], label.upper(), 9.5, col, bold=True)
gy, gh = 2.25, 3.95
group(s, ML, gy, 2.4, gh, "Browser", ACCENT)
group(s, ML+2.7, gy, 4.25, gh, "Remix server  /api/storefront/*")
group(s, ML+7.25, gy, CW-7.25, 1.95, "Supabase / Postgres", POS)
group(s, ML+7.25, gy+2.1, CW-7.25, 1.85, "AI providers", ACCENT)
# browser
chat = labelbox(s, ML+0.15, gy+0.35, 2.1, 0.75, "gleame-chat.js", sub="widget · state machine", anchor=MSO_ANCHOR.MIDDLE)
labelbox(s, ML+0.15, gy+1.3, 2.1, 0.7, "gleame-camera.js", sub="desktop webcam", anchor=MSO_ANCHOR.MIDDLE)
# server (2 cols)
sx1, sx2 = ML+2.85, ML+4.95; sbw = 1.9
cfg  = labelbox(s, sx1, gy+0.35, sbw, 0.68, "chat-config", sub="copy · axes", anchor=MSO_ANCHOR.MIDDLE)
rec  = labelbox(s, sx2, gy+0.35, sbw, 0.68, "recommend", sub="the brain", anchor=MSO_ANCHOR.MIDDLE)
txi  = labelbox(s, sx1, gy+1.25, sbw, 0.68, "transform-image", sub="per-product", anchor=MSO_ANCHOR.MIDDLE)
trk  = labelbox(s, sx2, gy+1.25, sbw, 0.68, "track-event", sub="funnel", anchor=MSO_ANCHOR.MIDDLE)
whk  = labelbox(s, ML+3.85, gy+2.95, sbw, 0.68, "webhooks/orders", sub="attribution", anchor=MSO_ANCHOR.MIDDLE)
# data
dx = ML+7.4; dw = CW-7.55
rules = labelbox(s, dx, gy+0.3, dw, 0.48, "recommendation_rules", anchor=MSO_ANCHOR.MIDDLE, lsize=11.5)
aev   = labelbox(s, dx, gy+0.85, dw, 0.48, "analytics_events", anchor=MSO_ANCHOR.MIDDLE, lsize=11.5)
word  = labelbox(s, dx, gy+1.4, dw, 0.48, "widget_orders", anchor=MSO_ANCHOR.MIDDLE, lsize=11.5)
# ai
gem = labelbox(s, dx, gy+2.4, dw, 0.55, "Google Gemini", sub="image-gen + vision · primary", anchor=MSO_ANCHOR.MIDDLE, lsize=12)
oai = labelbox(s, dx, gy+3.05, dw, 0.5, "OpenAI gpt-image", sub="fallback", anchor=MSO_ANCHOR.MIDDLE, lsize=12)
# edges
arrow(s, ML+2.25, gy+0.7, sx1, gy+0.69)                       # chat->config
arrow(s, ML+2.25, gy+0.85, sx1, gy+1.5)                       # chat->transform
arrow(s, sx2+sbw, gy+0.6, dx, gy+0.55, color=SUB)            # recommend->rules
arrow(s, sx2+sbw-0.2, gy+1.0, dx, gy+2.6, color=ACCENT, w=1.5)# recommend->gemini (AI)
arrow(s, sx2+sbw, gy+1.55, dx, gy+1.08)                       # track->analytics
arrow(s, dx+dw/2, gy+1.33, dx+dw/2, gy+1.4, color=ACCENT, dashed=True, w=1.25)  # cart_token aev->word
arrow(s, ML+5.75, gy+3.1, dx, gy+1.55, color=SUB)           # webhooks->orders
line(s, dx+dw/2, gy+2.95, dx+dw/2, gy+3.05, color=SUB, dashed=True)  # gemini->openai fallback
# legend
_, tf = textbox(s, ML, 6.35, CW, 0.35); p = tf.paragraphs[0]
_run(p, "——  ", 12, ACCENT, bold=True); _run(p, "AI call        ", 11.5, SUB)
_run(p, "----  ", 12, ACCENT, bold=True); _run(p, "cart_token join (analytics_events → widget_orders)", 11.5, SUB)
notes(s, "Walk left to right: browser widget, server endpoints, data and AI. Emphasize no iframe — same-origin gives direct cart access, speed, and trust. Legend: solid blue is an AI call; dashed is the cart_token join used for attribution.")

# 8 LIFECYCLE
s = slide(); header(s, "Request lifecycle", "A single recommendation request", idx())
left = ["1 · Boot — fetch chat-config + recommendation-config",
        "2 · Consult — answers accumulate into a criteria object",
        "3 · Photo — selfie posted to /chat-recommend",
        "4 · Gauntlet — verify shop, billing, rate-limit, 5 MB cap",
        "5 · Classify — Gemini reads the selfie into axes"]
right = ["6 · Match — JSONB-equality lookup, ordered by rank",
         "7 · Render — top-N try-ons generated in parallel",
         "8 · Backfill — a failed generation is replaced",
         "9 · Resolve — batched Admin GraphQL for product handles",
         "10 · Return — cards rendered to the shopper"]
bullets(s, ML, 2.05, (CW-0.6)/2, 4.2, left, size=14, gap=11)
bullets(s, ML+(CW-0.6)/2+0.6, 2.05, (CW-0.6)/2, 4.2, right, size=14, gap=11)
b = box(s, ML, 6.05, CW, 0.85, fill=PANEL, border=HAIR); b.text_frame.vertical_anchor=MSO_ANCHOR.MIDDLE
_run(b.text_frame.paragraphs[0], "Presented to the shopper as a calm three-step consultation; underneath, a security pass, a classification, a deterministic lookup, and N parallel generations with backfill.", 13, INK)
notes(s, "Narrate one shopper. The roughly twenty seconds is a UX estimate I still need to instrument end-to-end — say so. It becomes the central finding later.")

# 9 COLD START
s = slide(); header(s, "Recommendations", "Recommendation under cold-start", idx())
card(s, ML, 2.15, (CW-0.5)/2, 2.4, "The constraint",
     "A newly installed merchant has no click or purchase history. A learned recommender has nothing to train on at install.", body_size=15)
card(s, ML+(CW-0.5)/2+0.5, 2.15, (CW-0.5)/2, 2.4, "The approach",
     "Elicit preferences with a couple of questions; infer the rest from the photo; map the result to a merchant-authored rules matrix.", accent=True, body_size=15)
b = box(s, ML, 4.85, CW, 1.05, fill=PANEL, border=HAIR); b.text_frame.vertical_anchor=MSO_ANCHOR.MIDDLE
_run(b.text_frame.paragraphs[0], "Deterministic, debuggable, and merchant-controlled — it works at zero users, and is the substrate to layer learning onto once volume justifies it.", 14.5, INK)
notes(s, "This is the strongest design insight — slow down. Normal recommenders need data you do not have on day one. Inverting to elicit-and-infer is the answer; learning is a later upgrade, not a prerequisite.")

# 10 THE MATCH
s = slide(); header(s, "Recommendations", "From selfie to ranked products", idx())
c1 = card(s, ML, 2.15, (CW-0.5)/2, 3.6, "Photo classifier — a labeler, not a generator", "", body_size=13)
bullets(s, ML+0.2, 2.75, (CW-0.5)/2-0.4, 2.9,
        ["gemini-2.5-flash used as a vision classifier",
         "output constrained to an enum per axis",
         "temperature 0 — deterministic; cannot emit an unknown value",
         "768 px · 12 s timeout · errors degrade to no criteria"], size=13, gap=8)
c2 = box(s, ML+(CW-0.5)/2+0.5, 2.15, (CW-0.5)/2, 3.6, fill=WHITE, border=HAIR)
tf = c2.text_frame; tf.vertical_anchor=MSO_ANCHOR.TOP
_run(tf.paragraphs[0], "THE MATCH — STRICT JSONB EQUALITY", 10.5, ACCENT, bold=True)
code = tf.add_paragraph(); code.space_before=Pt(8)
_run(code, "SELECT variant_id, product_id, rank\nFROM recommendation_rules\nWHERE shop_id = $1\n  AND criteria = $2::jsonb\nORDER BY rank;", 12, INK, name=MONO)
note = tf.add_paragraph(); note.space_before=Pt(10)
_run(note, "A rule matches or it does not — no fuzzy scoring. One indexed hit; no match falls back to a shuffled, product-diverse set.", 12.5, SUB)
notes(s, "Two parts: a classifier that is a labeler (enum + temperature 0 means it cannot return a value the matrix does not understand) and a deterministic JSONB-equality lookup. Say plainly: no ML scoring — that is deliberate.")

# 11 IMAGE PIPELINE
s = slide(); header(s, "Image generation", "Personalized try-on generation", idx())
rows = [["Model", "Role"],
        [("gemini-3-pro-image", INK, True), "Variant configs — precise makeup"],
        [("gemini-2.5-flash-image", INK, True), "Standard — fast and low cost"],
        [("gemini-3.1-flash (2K)", INK, True), "High-detail input"],
        [("gpt-image-1.5 (OpenAI)", INK, True), "Fallback and refusal handling"]]
row_table(s, ML, 2.15, 6.0, rows, colx=[0.0, 0.52], fs=13, rowh=0.55)
# flow on right
fx = ML+6.6; fw = CW-6.6; fy=2.2
flow = ["Build prompt + fetch references (SSRF-safe)",
        "Compress · resize · HEIC→JPEG · EXIF-rotate",
        "Gemini · 2 retries with backoff",
        "On failure → OpenAI · safety-framed prompt"]
for i_, st in enumerate(flow):
    labelbox(s, fx, fy, fw, 0.62, st, fill=(ACCENT_SOFT if i_==3 else PANEL),
             border=(ACCENT if i_==3 else HAIR), align=PP_ALIGN.LEFT, lsize=12, lbold=False, anchor=MSO_ANCHOR.MIDDLE)
    if i_ < 3: arrow(s, fx+fw/2, fy+0.62, fx+fw/2, fy+0.78, color=SUB)
    fy += 0.8
b = box(s, ML, 5.75, CW, 1.0, fill=PANEL, border=HAIR); b.text_frame.vertical_anchor=MSO_ANCHOR.MIDDLE
_run(b.text_frame.paragraphs[0], "Top-N generations run in parallel, so wall-clock ≈ one image; backfill guarantees a full result set. Outputs are not cached — each is a function of one shopper's photo.", 13.5, INK)
notes(s, "Model tiering is a cost and quality decision, not only latency. The fallback chain is reliability engineering — single-provider face editing is fragile. The safety-framed prompt is how the OpenAI path avoids refusals.")

# 12 TRADEOFFS
s = slide(); header(s, "Key constraint", "Latency, cost, and quality", idx())
cols = [("Latency", "Parallel generation, input downscaling, capped retries, and a consultation UI that makes the wait purposeful."),
        ("Cost", "Flash by default and pro only where precision pays; generate only what is shown; rate limits as a circuit-breaker."),
        ("Quality", "Real product reference images, a deep-skin visibility strategy, and the pro model for makeup precision.")]
cw3 = (CW-0.8)/3
for i_, (t_, b_) in enumerate(cols):
    card(s, ML+i_*(cw3+0.4), 2.15, cw3, 2.7, t_, b_, body_size=13.5)
b = box(s, ML, 5.2, CW, 1.0, fill=PANEL, border=HAIR); b.text_frame.vertical_anchor=MSO_ANCHOR.MIDDLE
p = b.text_frame.paragraphs[0]
_run(p, "These three pull against each other; nearly every decision is a response. Even optimized, an estimated ~20 s wait — mostly on mobile — precedes value. ", 14, INK)
_run(p, "That estimate is the finding the data later confirms.", 14, INK, bold=True)
notes(s, "The central engineering tension. Name it explicitly: even optimized, value arrives after a perceived ~20 s, mostly on mobile — and that is exactly the funnel leak. Owning it reads as judgment.")

# 13 UNIT ECONOMICS
s = slide(); header(s, "Unit economics", "Flat revenue, variable cost", idx())
card(s, ML, 2.15, (CW-0.5)/2, 2.5, "Revenue — flat",
     "Session-based subscription tiers: Free $0 · Starter $30 · Launch $149 · Growth $399. The merchant pays the same regardless of try-on volume.", body_size=14)
c = box(s, ML+(CW-0.5)/2+0.5, 2.15, (CW-0.5)/2, 2.5, fill=WHITE, border=HAIR)
tf=c.text_frame; _run(tf.paragraphs[0], "COST — VARIABLE, PER SESSION", 10.5, ACCENT, bold=True)
cp=tf.add_paragraph(); cp.space_before=Pt(8); _run(cp, "cost = N × image-gen + 1 × classification", 14, INK, name=MONO)
cp2=tf.add_paragraph(); cp2.space_before=Pt(8); _run(cp2, "N = 3 by default (configurable 1–5). Scales with every completed flow. Per-image price comes from the provider rate card — not asserted from memory.", 13, SUB)
b = box(s, ML, 4.85, CW, 1.05, fill=PANEL, border=HAIR); b.text_frame.vertical_anchor=MSO_ANCHOR.MIDDLE
_run(b.text_frame.paragraphs[0], "Flat revenue against variable cost is why model tiering, generate-only-what-is-shown, and rate limits are margin decisions — not only performance ones.", 14.5, INK)
notes(s, "A consumer founder will ask unit economics. Lead with the structure. Do not quote a dollar figure that is not in the codebase; give the formula and, only if pushed, a pre-measured rate-card number.")

# 14 SECURITY
s = slide(); header(s, "Security", "Multi-tenant request hardening", idx())
rows = [["Control", "Purpose"],
        ["Shop-domain verification", "Prevents cross-shop data access and rate-limit spoofing"],
        ["Billing gate", "Blocks unpaid usage (active / trial / grace / grandfathered)"],
        ["Tiered rate limiting", "20/min · 100/hr per IP · 1000/hr per shop · 10/min on recommend"],
        ["Input validation", "5 MB cap, MIME and extension checks"],
        ["SSRF-safe fetch", "Guards merchant-supplied reference-image URLs"],
        ["Deterministic prompts", "No free-text reaches the model — no injection surface"]]
row_table(s, ML, 2.1, CW, rows, colx=[0.0, 0.34], fs=13, rowh=0.55)
_, tf = textbox(s, ML, 6.25, CW, 0.5)
_run(tf.paragraphs[0], "Honest limitation: the rate limiter is in-memory today — correct for current scale, and Redis-backed to scale horizontally.", 13, SUB)
notes(s, "Every storefront request crosses this boundary before doing expensive work. Volunteer the in-memory rate-limiter caveat — owning limits reads better than hiding them.")

# 15 ATTRIBUTION
s = slide(); header(s, "Measurement", "Funnel analytics and purchase attribution", idx())
flow = [("trackEvent", "client · fire-and-forget"), ("analytics_events", "+ cart_token"),
        ("orders/create", "Shopify webhook"), ("widget_orders", "+ journey data"),
        ("get_conversion_stats()", "attributed revenue")]
bw_ = 2.05; gap=(CW-5*bw_)/4; ty=2.8
xs=[]
for i_, (t_, sub_) in enumerate(flow):
    x = ML + i_*(bw_+gap); xs.append(x)
    accent = (i_==4)
    labelbox(s, x, ty, bw_, 0.95, t_, sub=sub_, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER,
             fill=(ACCENT_SOFT if accent else WHITE), border=(ACCENT if accent else HAIR), lsize=12)
    if i_ < 4: arrow(s, x+bw_+0.03, ty+0.47, x+bw_+gap-0.03, ty+0.47, color=SUB)
_, tf = textbox(s, ML, 4.2, CW, 0.8)
_run(tf.paragraphs[0], "The cart token is the only reliable session-to-order join on Shopify. Counts are event volume, not unique sessions (no session ID yet) — directional, and sufficient for relative funnel health.", 14, SUB)
notes(s, "Tracking must never slow the shopper path, so it is fire-and-forget; real attribution is reconciled server-side via the cart token and the orders webhook. Flag the events-not-sessions limitation yourself.")

# 16 RESULTS FUNNEL
s = slide(); header(s, "Results", "Funnel: one step dominates drop-off", idx())
data = [("Assistant opened", 245, None, INK),
        ("Consultation started", 235, "96%", POS),
        ("Photo uploaded", 79, "34% — 66% drop", NEG),
        ("Recommendations shown", 78, "99%", POS),
        ("Product clicked", 38, "49%", POS)]
bx = ML+3.0; maxw = 6.4; scale = maxw/245.0; yy=2.2; bh=0.42
for name_, val, pct, col in data:
    _, tf = textbox(s, ML, yy-0.03, 2.9, bh+0.1, anchor=MSO_ANCHOR.MIDDLE)
    _run(tf.paragraphs[0], name_, 12.5, INK)
    w = max(val*scale, 0.2)
    box(s, bx, yy, w, bh, fill=(NEG if col==NEG else ACCENT), border=None)
    _, tf = textbox(s, bx+w+0.1, yy-0.03, 1.0, bh+0.1, anchor=MSO_ANCHOR.MIDDLE)
    _run(tf.paragraphs[0], str(val), 13, INK, bold=True)
    if pct:
        _, tf = textbox(s, bx+w+0.7, yy-0.03, 3.2, bh+0.1, anchor=MSO_ANCHOR.MIDDLE)
        _run(tf.paragraphs[0], pct, 12, col, bold=(col==NEG))
    yy += 0.66
b = box(s, ML, 5.7, CW, 1.0, fill=PANEL, border=NEG, bw=1.25); b.text_frame.vertical_anchor=MSO_ANCHOR.MIDDLE
p=b.text_frame.paragraphs[0]
_run(p, "Diagnosis:  ", 14.5, NEG, bold=True)
_run(p, "every stage after a photo converts well. The business is gated on a single step — getting a mobile shopper to upload a photo.", 14.5, INK)
_, tf = textbox(s, ML, 6.78, CW, 0.3)
_run(tf.paragraphs[0], "Live on a production Shopify store · last 7 days · event volume, not unique sessions.", 10, SUB)
notes(s, "Real production data — most candidates have none. The funnel is healthy everywhere except the photo step, a 66% drop. Everything downstream converts. The whole business is gated on that one mobile step.")

# 17 RESULTS A/B
s = slide(); header(s, "Results", "A/B test: negative in aggregate, positive on desktop", idx())
rows = [["Segment", "Conversion", "AOV", "Revenue / visitor"],
        ["Blended", ("▼ 3.0% vs 3.7%", NEG, True), ("▲ +7.5%", POS, True), ("▼ down", NEG, False)],
        ["Desktop", ("▲ up", POS, True), ("▲ up", POS, True), ("▲ up", POS, True)],
        ["Mobile  (majority of traffic)", ("▼ down", NEG, True), ("~ flat", SUB, False), ("▼ down", NEG, True)]]
row_table(s, ML, 2.2, CW, rows, colx=[0.0, 0.34, 0.56, 0.76], fs=14, rowh=0.62)
b = box(s, ML, 5.1, CW, 1.1, fill=PANEL, border=HAIR); b.text_frame.vertical_anchor=MSO_ANCHOR.MIDDLE
p=b.text_frame.paragraphs[0]
_run(p, "Higher AOV with lower conversion is signal, not noise: ", 14, INK, bold=True)
_run(p, "Gleame attracts higher-basket buyers but loses marginal mobile converters at the photo step.", 14, INK)
_, tf = textbox(s, ML, 6.3, CW, 0.4)
_run(tf.paragraphs[0], "Shopify-native experiment. Sample ~60–74 conversions per arm — directional, not yet statistically significant.", 11, SUB)
notes(s, "Say clearly the A/B is Shopify-native — you built the product and instrumentation, not the experiment framework. Blended is negative, but the device split is the story: desktop wins, mobile loses, mobile dominates traffic. Volunteer the small-sample caveat.")

# 18 DIAGNOSIS
s = slide(); header(s, "Diagnosis", "Two signals, one root cause", idx())
card(s, ML, 2.15, (CW-0.5)/2, 2.2, "Signal 1 — the funnel", "A 66% drop at the mobile photo step; everything downstream converts.", body_size=15)
card(s, ML+(CW-0.5)/2+0.5, 2.15, (CW-0.5)/2, 2.2, "Signal 2 — the A/B test", "Desktop (frictionless photo) wins; mobile loses and drags the blend negative.", body_size=15)
b = box(s, ML, 4.6, CW, 1.5, fill=ACCENT_SOFT, border=ACCENT, bw=1.25); b.text_frame.vertical_anchor=MSO_ANCHOR.MIDDLE
p=b.text_frame.paragraphs[0]
_run(p, "The model works; the mobile photo step is the bottleneck. ", 16, INK, bold=True)
_run(p, "This is a product fix, not a model fix — and the desktop result already proves the concept when shoppers clear that step.", 16, INK)
notes(s, "The conclusion. Two independent sources triangulate to one cause: mobile photo friction. Deliver it plainly and without theatrics. The next move is product, not more AI.")

# 19 ROADMAP
s = slide(); header(s, "Roadmap", "What I would build next", idx())
rm = [("Reduce mobile photo friction", "Stream one fast try-on first, backfill the rest; tighten camera UX and framing copy."),
      ("Session identifiers", "Honest unique-user funnels and true per-session conversion."),
      ("Layer learning on the matrix", "Re-rank merchant rules by observed clicks and conversions."),
      ("Measure returns", "Validate the second-order thesis: accuracy reduces returns."),
      ("Scale-out infrastructure", "Redis-backed rate limiter; a queue with streamed results."),
      ("Targeted caching", "Cache classifications and reference compressions, not final outputs.")]
cw2=(CW-0.5)/2; ch=1.35
for i_, (t_, b_) in enumerate(rm):
    r,c = divmod(i_,2)
    card(s, ML+c*(cw2+0.5), 2.05+r*(ch+0.3), cw2, ch, f"{i_+1} · {t_}", b_, body_size=12.5)
notes(s, "Shows you think in roadmaps. Item one directly attacks the diagnosed leak. Returns measurement is the unmeasured second-order value.")

# 20 SUMMARY
s = slide(); header(s, "Summary", "Summary", idx())
bullets(s, ML, 2.1, CW, 3.2,
        [("A complete CRO system, ", "not a model wrapper: recommender, image pipeline, prompt layer, security boundary, and attribution."),
         ("Built for cold-start: ", "elicit and infer preferences, then map to a deterministic, merchant-controlled matrix."),
         ("Instrumented and A/B-tested in production, ", "then diagnosed honestly."),
         ("The result is actionable: ", "desktop validates the concept; one mobile-friction fix is the next unlock.")],
        size=15.5, gap=13)
b = box(s, ML, 5.6, CW, 0.95, fill=PANEL, border=HAIR); b.text_frame.vertical_anchor=MSO_ANCHOR.MIDDLE
_run(b.text_frame.paragraphs[0], "Thank you. I'm happy to go deeper on any subsystem.", 15, INK, bold=True)
notes(s, "Close on the reframe and the actionable conclusion. Invite hard questions — you have the answer bank cold. If asked unit economics, give the formula and a pre-measured number, never an improvised dollar figure.")

prs.save("/Users/charles/glimpse-app/gleame-assistant-prep/Gleame_Technical_Deep_Dive.pptx")
print(f"Saved {n} content slides + title =", len(prs.slides.__iter__.__self__._sldIdLst), "total slides")
